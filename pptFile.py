from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx import Presentation
import tempfile
import os
import copy
import pickle
import re

def process_PPTfile(uploadedFile):
    # Load the original PPTX presentation
    pptx_file_path = save_file(uploadedFile,'uploadFile.pptx')
    prs = Presentation(pptx_file_path)
    # Process each slide in the presentation
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    newparagraph = copy.deepcopy(paragraph)
                    seperate_important_words(paragraph,newparagraph)

    modified_pptx_file = 'modified_presentation.pptx'
    # Save the modified presentation 
    processed_ppt_path = save_file(prs,modified_pptx_file)
    print(processed_ppt_path)
    return modified_pptx_file

def seperate_important_words(paragraph,newparagraph):
    paragraph.clear()
    for run in newparagraph.runs:
        words = run.text.split()
        for word in words:
            #seperates all the words in their own seperate runs
            cleanword = clean_word(word)
            if (important_word_satisfied(cleanword)):
                # Create a new run with modified formatting
                new_run = paragraph.add_run()
                new_run.text = word + " "
                new_run.font.size = run.font.size  # Modify the font size as needed
                add_hyperlink(new_run,f'https://en.wikipedia.org/wiki/{word}',new_run.text)
            else:
                # Create a new run for non-important words (preserving existing formatting)
                new_run = paragraph.add_run()
                new_run.text = word + " "
                new_run.font.size = run.font.size  # Preserve font size

#condition
def clean_word(word):
    cleaned_word = re.sub(r'[^a-zA-Z\'’’]','',word) #maybe just don't even consider the words with apostrophies for now
    return cleaned_word.lower()

def get_wordset_frmFile(filepath):
    with open(filepath,'rb') as file:
        word_set = pickle.load(file)
    return word_set
def important_word_satisfied(word):
    word_set = get_wordset_frmFile('/Users/venkataneti/Desktop/pptApp/uploads/wordset')
    if(word not in word_set):
        with open('words','a') as file:
            file.write(f"{word}\n")
        return True
    return False

# Function to add hyperlinks to text runs
def add_hyperlink(run, url, text):
    run.hyperlink.address = url
    run.text = text

def save_file(presentation,name):
    # Create a temporary directory to store the processed file
    temp_dir = '/Users/venkataneti/Desktop/pptApp/uploads'
    file_path = os.path.join(temp_dir, name)
    
    # Save the processed presentation
    presentation.save(file_path)
    
    return file_path

